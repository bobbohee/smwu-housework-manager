#!/usr/bin/env python3
"""우리집 살림 매니저 제안서 PPTX 생성 (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "우리집-살림-매니저-제안서.pptx")

KFONT = "Apple SD Gothic Neo"
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x5C, 0x9A)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF5, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_font(run, size, bold=False, color=DARK, font=KFONT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def fill_rect(slide, l, t, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def header(slide, title, kicker=None):
    fill_rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    tf = textbox(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.85),
                 anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    if kicker:
        r = p.add_run(); r.text = kicker + "  "
        set_font(r, 14, bold=True, color=RGBColor(0x9F, 0xC0, 0xE8))
    r = p.add_run(); r.text = title
    set_font(r, 26, bold=True, color=WHITE)


def bullets(tf, items, size=16, gap=6):
    """items: list of (text, level, bold)."""
    first = True
    for it in items:
        text, level, bold = (it + (False,))[:3] if isinstance(it, tuple) else (it, 0, False)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.space_after = Pt(gap)
        prefix = "•  " if level == 0 else "–  "
        r = p.add_run(); r.text = prefix + text
        set_font(r, size - level * 2, bold=bold, color=DARK if level == 0 else GRAY)


def add_table(slide, l, t, w, rows, col_widths, header_row=True, fontsize=12):
    nrows = len(rows); ncols = len(rows[0])
    h = Inches(0.0)
    gtbl = slide.shapes.add_table(nrows, ncols, l, t, w, Inches(0.4 * nrows))
    tbl = gtbl.table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            if header_row and ri == 0:
                set_font(r, fontsize, bold=True, color=WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                set_font(r, fontsize, color=DARK)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
    return gtbl


# ───────────────────────── 1. 표지 ─────────────────────────
s = add_slide()
fill_rect(s, 0, 0, SW, SH, NAVY)
fill_rect(s, 0, Inches(2.7), SW, Inches(2.1), RGBColor(0x16, 0x2B, 0x47))
tf = textbox(s, Inches(1), Inches(2.85), Inches(11.3), Inches(1.9), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "우리집 살림 매니저"
set_font(r, 48, bold=True, color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "함께 사는 사람들을 위한 집안일 당번·기록 웹앱"
set_font(r, 20, color=RGBColor(0xBF, 0xD4, 0xEE))
tf3 = textbox(s, Inches(1), Inches(6.4), Inches(11.3), Inches(0.6))
p = tf3.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "웹시스템설계 기말 프로젝트 제안서  ·  2026-05-24"
set_font(r, 14, color=RGBColor(0x9F, 0xC0, 0xE8))

# ───────────────────────── 1-2. 목차 ─────────────────────────
s = add_slide(); header(s, "목차", "CONTENTS")
TOC = [
    ("01", "프로젝트 개요", "서비스 정의·대상·핵심 가치"),
    ("02", "문제 정의 / 대상 사용자", "공동 거주 분담 갈등, 페르소나"),
    ("03", "시장 및 유사 서비스 분석", "경쟁 비교, 차별성(Positioning)"),
    ("04", "요구사항 분석", "기능/비기능, MoSCoW, 필수요건 매핑"),
    ("05", "주요 기능 설명", "그룹·집안일·완료·비활성화·꽝뽑기·캘린더·통계"),
    ("06", "화면 설계", "화면설계서 13종 (와이어프레임 + Description)"),
    ("07", "기술 스택 및 시스템 구조", "기술 선택, 3-Tier 구조, 데이터 모델"),
    ("08", "개발 일정", "3주 집중 일정, 중간 점검"),
]
ty = Inches(1.5)
for num, title, desc in TOC:
    fill_rect(s, Inches(0.7), ty, Inches(0.62), Inches(0.56), BLUE)
    ntf = textbox(s, Inches(0.7), ty, Inches(0.62), Inches(0.56), anchor=MSO_ANCHOR.MIDDLE)
    p = ntf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; set_font(r, 16, bold=True, color=WHITE)
    ttf = textbox(s, Inches(1.55), ty, Inches(11.2), Inches(0.56), anchor=MSO_ANCHOR.MIDDLE)
    p = ttf.paragraphs[0]
    r = p.add_run(); r.text = title + "   "; set_font(r, 17, bold=True, color=NAVY)
    r = p.add_run(); r.text = desc; set_font(r, 12, color=GRAY)
    ty += Inches(0.66)
tf = textbox(s, Inches(1.55), ty + Inches(0.04), Inches(11.2), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "부록 · AI 활용 계획"; set_font(r, 13, bold=True, color=GRAY)

# ───────────────────────── 2. 프로젝트 개요 ─────────────────────────
s = add_slide(); header(s, "프로젝트 개요", "01")
add_table(s, Inches(0.6), Inches(1.5), Inches(12.1), [
    ["항목", "내용"],
    ["서비스명", "우리집 살림 매니저"],
    ["한 줄 정의", "공동 거주 구성원의 집안일 당번·완료 기록을 관리하는 웹앱"],
    ["대상 사용자", "함께 거주하는 2~N명 (룸메이트·가족·셰어하우스 등)"],
    ["핵심 가치", "“지금 누구 차례인가”를 다툼 없이 한눈에 + 기여도를 데이터로 가시화"],
    ["플랫폼", "모바일 우선 반응형 웹 (PC 사이드바 / 모바일 하단 탭)"],
], [Inches(2.6), Inches(9.5)], fontsize=14)
tf = textbox(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(2.2))
bullets(tf, [
    "기존: 가족 캘린더 앱에 “OO 설거지”처럼 수동 기록 → 순번 자동화·집계가 안 됨",
    "전용 도구로 대체: 순번 자동 진행 · 완료 기록 · 통계 · 랜덤 배정",
    "특정 가정에 고정되지 않은 범용 서비스 (멤버·집안일·색상을 그룹이 직접 정의)",
], size=16)

# ───────────────────────── 3. 문제 정의 ─────────────────────────
s = add_slide(); header(s, "문제 정의", "02")
tf = textbox(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.2))
bullets(tf, [
    ("공동 거주의 집안일 분담 = 사소하지만 반복되는 갈등의 원인", 0, True),
    ("“누가 할 차례인지” 불명확 — 기억·구두 약속 의존, 책임 소재 흐려짐", 1),
    ("수동 기록의 한계 — 순번 자동 진행·집계 불가, 누락·중복 잦음", 1),
    ("기여도 불투명 — 월 기여 수치 없어 “나만 한다” 불만 누적", 1),
    ("돌발 집안일 배정의 비효율 — 공정한 배정 수단 부재", 1),
], size=16)
fill_rect(s, Inches(0.6), Inches(4.9), Inches(12.1), Inches(1.9), LIGHT)
tf = textbox(s, Inches(0.9), Inches(5.05), Inches(11.5), Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
bullets(tf, [
    ("핵심 질문", 0, True),
    ("① 어떤 불편함? 차례·기여도 불투명   ② 왜? 자동 기록·집계 수단 부재", 1),
    ("③ 해결 시 가치? 분담 갈등 감소 + 공정성 가시화", 1),
], size=15)

# ───────────────────────── 4. 대상 사용자 ─────────────────────────
s = add_slide(); header(s, "대상 사용자 (Persona)", "02")
add_table(s, Inches(0.6), Inches(1.6), Inches(12.1), [
    ["페르소나", "상황", "니즈"],
    ["공유 룸메이트 (20대)", "셰어하우스 3인, 생활 패턴 제각각", "차례 자동 안내, 안 한 사람 확인"],
    ["가족 구성원 (자매 3인)", "기존 캘린더로 수기 관리", "수기 대체, 월별 기여도 통계"],
    ["방장 (그룹 관리자)", "분담 규칙·멤버 관리 책임", "멤버 초대/관리, 규칙 위반 기록 관리"],
], [Inches(3.2), Inches(4.7), Inches(4.2)], fontsize=14)

# ───────────────────────── 5. 시장/유사 서비스 ─────────────────────────
s = add_slide(); header(s, "시장 및 유사 서비스 분석", "03")
add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), [
    ["서비스", "분류", "장점", "한계"],
    ["Tody / 집안일 체크앱", "청소 관리 앱", "주기 기반 청소 알림", "개인용 중심, 공동 순번·비교 약함"],
    ["가족 공유 캘린더", "일반 캘린더", "익숙함, 일정 공유", "순번 자동·집계·랜덤 배정 없음"],
    ["카카오톡 공지", "메신저", "진입장벽 0", "기록 휘발, 통계·추적 불가"],
    ["종이 당번표", "아날로그", "즉시성", "원격 확인·데이터화 불가"],
], [Inches(2.9), Inches(2.0), Inches(3.4), Inches(3.8)], fontsize=13)

# ───────────────────────── 6. 차별성 ─────────────────────────
s = add_slide(); header(s, "차별성 (Positioning)", "03")
tf = textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5))
bullets(tf, [
    ("“공동 거주 전용 + 순번 자동화 + 기여도 통계”의 교집합", 0, True),
    ("순번제 자동 진행 — 완료 1회 클릭에 다음 사람으로 자동 이행", 1),
    ("고정제 + 순번제 혼합 — 요일/주기 고정 담당과 순번 돌리기 동시 운용", 1),
    ("기여도 통계 + 캘린더 — 월별 멤버별/집안일별 완료를 수치·차트로 가시화", 1),
    ("꽝뽑기(랜덤 배정) — 돌발 집안일을 공정하게 배정", 1),
    ("규칙 위반 관리 — 완료 기록 비활성화(사유) + 순번 복원으로 공정성 담보", 1),
], size=17)

# ───────────────────────── 7. 요구사항 (기능) ─────────────────────────
s = add_slide(); header(s, "요구사항 분석 · 기능 요구사항", "04")
add_table(s, Inches(0.6), Inches(1.35), Inches(12.1), [
    ["ID", "기능", "설명"],
    ["FR-01", "인증", "이메일/비밀번호 회원가입·로그인 (Firebase Auth)"],
    ["FR-02", "그룹 관리", "그룹 생성, 초대 코드, 강퇴(방장), 방장 위임, 다중 그룹 전환"],
    ["FR-03", "집안일 CRUD", "프리셋 + 자유 추가/수정/삭제, 모드·색상·규칙 설정"],
    ["FR-04", "순번 완료 체크", "현재 차례 멤버 완료 → 기록 저장 + 차례 자동 이행"],
    ["FR-05", "완료 내역 비활성화", "방장이 사유와 함께 비활성화 + 순번 복원"],
    ["FR-06", "꽝뽑기", "참여 멤버 중 N명 랜덤 배정, 결과 기록"],
    ["FR-07", "캘린더 뷰", "월별 그리드에 완료 기록 색상별 표시, 상세 팝업"],
    ["FR-08", "통계", "월별 멤버별/집안일별 완료 횟수 (표 + 차트)"],
], [Inches(1.2), Inches(2.7), Inches(8.2)], fontsize=13)

# ───────────────────────── 8. 요구사항 (비기능 + MoSCoW) ─────────────────────────
s = add_slide(); header(s, "비기능 요구사항 · MoSCoW", "04")
tf = textbox(s, Inches(0.6), Inches(1.35), Inches(5.9), Inches(5))
bullets(tf, [
    ("비기능 요구사항", 0, True),
    ("성능: 주요 화면 3초 내 로딩, 실시간 리스너 즉시 반영", 1),
    ("보안: Firebase Auth + Firestore Rules (그룹 멤버만 접근)", 1),
    ("반응형: 1280 / 768 / 375px 대응", 1),
    ("가용성: 정적 호스팅 + 클라이언트-DB 직접 통신", 1),
], size=14)
add_table(s, Inches(6.7), Inches(1.5), Inches(6.0), [
    ["등급", "기능"],
    ["Must", "인증, 그룹, 집안일 CRUD, 순번 완료, 캘린더"],
    ["Should", "통계, 완료 비활성화, 꽝뽑기"],
    ["Could", "다크/라이트, 초대 링크, .ics 내보내기"],
    ["Won't", "푸시 서버, 구글 캘린더 동기화, AI 추천"],
], [Inches(1.3), Inches(4.7)], fontsize=13)

# ───────────────────────── 9. 필수요건 매핑 ─────────────────────────
s = add_slide(); header(s, "강의 필수요건 충족 매핑", "04")
add_table(s, Inches(0.6), Inches(1.35), Inches(12.1), [
    ["필수요건", "충족 방식"],
    ["최소 10개 화면", "13개 화면 설계"],
    ["입력 폼 1개 이상", "회원가입 · 집안일 등록 · 비활성화 사유 · 꽝뽑기 설정"],
    ["CRUD 3개 이상", "집안일 C/R/U/D + 완료기록 Create/Update"],
    ["상태 관리 1개 이상", "로그인 상태 · 현재 그룹 선택 · 캘린더 월/필터"],
    ["데이터 저장 / 외부 API", "Firebase Firestore + Firebase Auth"],
    ["반응형 UI", "모바일 하단탭 ↔ PC 사이드바"],
    ["AI 활용 / 배포", "Claude로 설계·코드 생성(기록) / Vercel 배포"],
], [Inches(3.3), Inches(8.8)], fontsize=13)

# ───────────────────────── 10. 주요 기능 (1/2) — 그룹·집안일 ─────────────────────────
s = add_slide(); header(s, "주요 기능 설명 (1/2) · 그룹 · 집안일", "05")
tf = textbox(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.9))
bullets(tf, [
    ("그룹 관리 — 데이터 단위", 0, True),
    ("그룹 생성자 = 방장(ownerId). groups.memberUids[] ↔ users.groupIds[] 양방향 동기화로 일관성 유지", 1),
    ("초대 코드: 영문 대문자+숫자 6자, 혼동문자(0/O·1/I/L) 제외, 그룹당 1개·충돌 시 재생성, 기간/횟수 제한 없음", 1),
    ("방장 전용: 멤버 강퇴 · 방장 위임(ownerId 이전) · 완료 내역 비활성화", 1),
    ("다중 그룹: 한 계정이 여러 그룹 가입, 홈 상단 그룹명 탭 → 앱 전체 데이터 컨텍스트 전환", 1),
    ("집안일 — 프리셋 & 색상", 0, True),
    ("그룹 생성 직후 프리셋 제시 → 선택 시 chores에 ‘복사본’ 생성(참조 아님). 0개(빈 시작)도 허용, 이후 자유 추가/수정/삭제", 1),
    ("후보: 순번제 5(설거지·빨래·거실/방 청소·밥) · 고정제 3(음식물/일반쓰레기·화장실)", 1),
    ("10색 고정 팔레트에서 1색 선택 → 캘린더 점·홈 카드·통계 막대에 동일 색 적용", 1),
    ("집안일 — 순번제 / 고정제 모드", 0, True),
    ("순번제(rotation): chore별 참여 멤버·순서 독립 지정, 주기 없음(쌓이면 차례인 사람 수행)", 1),
    ("고정제(fixed): 요일 지정(weekly) 또는 N일 주기(interval, 시작일 기준). 완료 기록 없이 홈 ‘오늘 담당’ 알림만", 1),
    ("완료 규칙(rules): 집안일별 약속을 참고용 표시(앱이 강제 체크하지 않음)", 1),
], size=14, gap=4)

# ───────────────────────── 10b. 주요 기능 (2/2) — 완료·비활성화·꽝뽑기·캘린더·통계 ─────────────────────────
s = add_slide(); header(s, "주요 기능 설명 (2/2) · 완료 · 기록 · 통계", "05")
tf = textbox(s, Inches(0.6), Inches(1.3), Inches(12.1), Inches(5.9))
bullets(tf, [
    ("완료 체크 & 권한 (순번제 전용)", 0, True),
    ("완료 1회 클릭 → ① choreLog 기록 생성(type:rotation, active:true) ② currentTurnIndex=(idx+1)%N 차례 이행", 1),
    ("기본은 현재 차례 멤버만 활성. allowProxyComplete=true면 참여 멤버 누구나 대신 완료(부재 커버)", 1),
    ("귀속 분리: completedBy=차례 멤버(통계 귀속) / completedByActual=실제 누른 사람. 통계는 completedBy 기준", 1),
    ("완료 내역 비활성화 (방장)", 0, True),
    ("사유 입력(필수) 후 active:false로 보존(삭제X) — deactivatedBy/At/Reason 기록, 전원 확인 가능", 1),
    ("통계 제외 + 순번제는 currentTurnIndex 되돌려 차례 복원. 캘린더엔 취소선·딤으로 구분 표시", 1),
    ("꽝뽑기 (랜덤 배정)", 0, True),
    ("멤버 체크박스 참여/제외, 당첨 인원 조정(기본 1명) → 무작위 N명, 폭탄 연출. choreLog에 type:random 기록(순번 무관)", 1),
    ("캘린더 · 통계", 0, True),
    ("자체 캘린더(외부 비의존): completedAt 기준 월 그리드, 색상 점 표시, 클릭 시 상세 팝업(담당자/실제 완료/시간/유형/상태)", 1),
    ("통계: 이번 달 멤버별·집안일별 완료 횟수를 표 + 차트로, active=true 기록만 집계", 1),
], size=14, gap=4)

# ───────────────────────── 11. 화면 설계 표지 ─────────────────────────
s = add_slide(); header(s, "화면 설계 (화면설계서)", "06")
tf = textbox(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(4.5))
bullets(tf, [
    ("표준 양식: 화면 코드 / Page title / Screen Path + 와이어프레임 + Description", 0, True),
    ("캡처 좌측의 번호(①②③…) 배지 = 우측 Description 항목 번호와 1:1 매칭", 1),
    ("인터랙티브 프로토타입을 PC(데스크톱)·라이트 테마로 캡처 (좌측 사이드바 네비게이션)", 1),
    ("총 13개 화면 — 강의 필수요건(최소 10화면) 충족", 1),
    ("네비게이션: 하단 탭(모바일)·사이드바(PC) — 모든 기능 2~3클릭 내 도달 목표", 1),
    ("색상 팔레트(10색 고정, 집안일 구분용)", 0, True),
    ("#4A90D9  #E74C3C  #2ECC71  #F39C12  #9B59B6  #1ABC9C  #E91E8C  #34495E  #95A5A6  #F1C40F", 1),
], size=16)
# 색상 스와치 (시각 표시)
palette = ["4A90D9", "E74C3C", "2ECC71", "F39C12", "9B59B6",
           "1ABC9C", "E91E8C", "34495E", "95A5A6", "F1C40F"]
sw = Inches(1.0); gap = Inches(0.18); sx = Inches(0.7); sy = Inches(4.7)
for i, hexc in enumerate(palette):
    c = RGBColor(int(hexc[0:2], 16), int(hexc[2:4], 16), int(hexc[4:6], 16))
    fill_rect(s, sx + i * (sw + gap), sy, sw, Inches(0.9), c)

# 화면설계 슬라이드 (이미지 + Description)
SCREENS = [
    ("01-auth-login", "SCREEN-01 · 로그인 / 회원가입", "앱 진입 시 최초 화면", [
        "서비스 로고(🏠)·명칭을 중앙 배치해 앱 정체성을 즉시 인지시키는 브랜드 영역",
        "이메일 입력 — Firebase Auth 계정 식별자, @형식 검증 후 진행",
        "비밀번호 입력 — 마스킹(••) 처리, 최소 길이 검증",
        "로그인 버튼 — 인증 성공 시 홈으로 라우팅, 실패 시 오류 메시지 노출",
        "회원가입 전환 링크 — 이름·비밀번호 확인 필드가 추가된 가입 폼으로 토글",
    ]),
    ("02-home", "SCREEN-02 · 홈 (대시보드)", "로그인 후 / 하단탭 ‘홈’", [
        "상단 그룹 바 — 그룹명(▼) 탭 시 그룹 전환 오버레이 호출, 우측에 오늘 날짜 표시",
        "오늘의 고정 집안일 알림 — 고정제 담당자 안내(완료 버튼 없음)",
        "순번제 카드 그리드 — 집안일별 이모지·색상·현재 차례 담당자 표시(PC 다열/모바일 2열)",
        "완료 버튼 — 차례 멤버 1회 클릭 시 완료 기록 저장 + currentTurnIndex 자동 +1(다음 멤버 이행)",
    ]),
    ("03-group-overlay", "SCREEN-03 · 그룹 전환 오버레이", "홈 상단 그룹명 탭", [
        "오버레이 제목 — ‘그룹 선택’, 배경 딤 처리로 포커스 집중",
        "현재 그룹 — 활성 그룹을 ‘현재’ 배지로 강조해 컨텍스트 식별",
        "다른 가입 그룹 — 클릭 시 데이터 컨텍스트 전환(홈·캘린더·통계 모두 갱신)",
        "새 그룹 만들기 / 6자 초대 코드 합류 진입점",
    ]),
    ("04-calendar", "SCREEN-04 · 캘린더", "하단탭 ‘캘린더’", [
        "월 이동 헤더 — ◀▶로 이전·다음 달 탐색, 중앙에 연·월 표시",
        "월간 날짜 그리드 — 완료 기록을 집안일 색상 점으로 누적 표시(다건은 여러 점)",
        "선택일 완료 기록 — 색상·담당자·시간 나열, 비활성 건은 취소선·딤, 클릭 시 상세 팝업",
    ]),
    ("05-calendar-popup", "SCREEN-05 · 완료 상세 팝업", "캘린더 → 기록 클릭", [
        "팝업 헤더 — 대상 집안일명과 닫기(✕)",
        "상세 정보 — 집안일·담당자(차례)·실제 완료자·완료 시간·유형·상태 행 표시(담당자/실제 완료자 분리)",
        "비활성화 진입 — 방장에게만 노출, 규칙 위반 기록 처리 폼으로 이동",
    ]),
    ("06-calendar-deactivate", "SCREEN-06 · 완료 내역 비활성화", "완료 상세 → ‘비활성화’", [
        "비활성화 대상 — 처리할 완료 기록(집안일·담당자·시간)을 재확인",
        "사유 입력 — 필수 입력, 미입력 시 확정 불가(공정성·추적성 확보)",
        "처리 버튼 — 확정 시 통계 집계 제외 + 소진된 순번 차례 복원 / 취소 시 상세 복귀",
    ]),
    ("07-chores", "SCREEN-07 · 집안일 관리", "설정 → 집안일 관리", [
        "헤더 — 뒤로·제목과 우측 ‘+ 추가’(신규 등록 진입)",
        "프리셋 빠른 추가 — 기본 항목을 칩 클릭으로 즉시 복사 생성(참조 아닌 복사 → 독립 편집 가능)",
        "집안일 목록 — 색상 점·모드(순번/고정)·참여 인원과 함께 나열, 클릭 시 수정",
    ]),
    ("08-chore-edit", "SCREEN-08 · 집안일 추가 / 수정 (순번제)", "집안일 관리 → 추가/클릭", [
        "집안일 이름 — 자유 입력(프리셋 복사 후 변경 가능)",
        "색상 — 10색 팔레트 선택, 캘린더 점·홈 카드·통계 막대에 동일 색 적용",
        "모드 — 순번제/고정제 토글, 선택에 따라 하단 전용 설정 전환",
        "참여 멤버·순서 — 체크로 참여 지정, 드래그(☰)로 순번 순서 조정",
        "대신 완료 허용(allowProxyComplete) — 켜면 차례 부재 시 타 멤버 대신 완료(담당자/실제 완료자 분리 기록)",
        "완료 규칙(참고용) — 완료 기준 메모, 멤버 간 합의 용도",
    ]),
    ("08b-chore-edit-fixed", "SCREEN-08B · 집안일 편집 (고정제)", "집안일 편집 → 모드 ‘고정제’", [
        "집안일 이름 — 자유 입력 (공통)",
        "색상 — 10색 팔레트 선택 (공통)",
        "모드 — 고정제 선택 상태",
        "스케줄 방식 — ‘요일 지정’ / ‘주기 지정’ 택1",
        "담당자별 요일 — 멤버마다 담당 요일 토글(요일 지정), 주기 지정 시 N일 주기+시작일(예: 2주마다)",
        "완료 규칙(참고용) — 고정제는 완료 버튼 없이 ‘오늘 담당’ 알림만 제공, 규칙은 안내용",
    ]),
    ("09-random", "SCREEN-09 · 꽝뽑기 설정", "하단탭 ‘꽝뽑기’", [
        "참여 멤버 — 체크박스로 추첨 대상 포함/제외",
        "당첨 인원 — −/+로 뽑을 인원 조정(기본 1명), ‘N명 중 M명’ 안내",
        "꽝뽑기 버튼 — 클릭 시 참여자 중 무작위 배정 실행",
    ]),
    ("10-random-result", "SCREEN-10 · 꽝뽑기 결과", "꽝뽑기 → 실행", [
        "발표 라벨 — 결과 연출 안내",
        "결과 — 당첨자에 폭탄(💣) 강조, 미당첨은 딤 처리",
        "액션 — ‘다시하기’ 재추첨 / ‘확인’ 시 random 유형으로 완료 기록 저장",
    ]),
    ("11-stats", "SCREEN-11 · 통계", "하단탭 ‘통계’", [
        "월 이동 헤더 — 집계 대상 월 선택(◀▶)",
        "멤버별 완료 횟수 — 막대 차트 + 표로 기여도 가시화(활성 기록만 집계)",
        "집안일별 완료 횟수 — 항목별 누적 횟수를 가로 막대로 비교",
    ]),
    ("12-settings", "SCREEN-12 · 그룹 설정", "하단탭 ‘설정’", [
        "그룹 이름 — 인라인 수정",
        "초대 코드 — 6자 코드 확인·복사·공유(신규 멤버 합류용)",
        "멤버 목록 — 방장(👑)·멤버 구분, 방장에게 ‘방장 위임’·‘강퇴’ 노출",
        "집안일 관리 진입 — 집안일 CRUD 화면으로 이동",
        "그룹 나가기 — 본인 탈퇴(방장은 위임 후 가능)",
    ]),
]

for img, title, path, descs in SCREENS:
    s = add_slide()
    fill_rect(s, 0, 0, SW, Inches(1.0), NAVY)
    tf = textbox(s, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    set_font(r, 20, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = "Screen Path:  " + path
    set_font(r, 12, color=RGBColor(0xBF, 0xD4, 0xEE))
    # 이미지 (좌측), PC 가로 캡처 → 폭 기준 배치 + 높이 상한
    imgpath = os.path.join(ASSETS, img + ".png")
    iw, ih = Image.open(imgpath).size
    disp_w = Inches(7.2)
    disp_h = Emu(int(disp_w * ih / iw))
    max_h = Inches(5.95)
    if disp_h > max_h:
        disp_h = max_h
        disp_w = Emu(int(disp_h * iw / ih))
    left = Inches(0.55)
    top = Inches(1.35)
    s.shapes.add_picture(imgpath, left, top, width=disp_w, height=disp_h)
    # Description 박스 (우측)
    dx = left + disp_w + Inches(0.4)
    dw = SW - dx - Inches(0.4)
    fill_rect(s, dx, Inches(1.35), dw, Inches(0.5), BLUE)
    htf = textbox(s, dx + Inches(0.15), Inches(1.4), dw - Inches(0.3), Inches(0.4), anchor=MSO_ANCHOR.MIDDLE)
    r = htf.paragraphs[0].add_run(); r.text = "Description"
    set_font(r, 16, bold=True, color=WHITE)
    rows = [["#", "설명"]] + [[str(i + 1), d] for i, d in enumerate(descs)]
    add_table(s, dx, Inches(2.0), dw, rows, [Inches(0.5), dw - Inches(0.5)], fontsize=12)

# ───────────────────────── 12. 기술 스택 ─────────────────────────
s = add_slide(); header(s, "기술 스택", "07")
add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), [
    ["영역", "기술", "선택 이유"],
    ["프론트엔드", "React + Next.js (App Router)", "컴포넌트 기반·라우팅, 상태 변화 중심 SPA"],
    ["스타일", "Tailwind CSS", "빠른 반응형 UI 구성"],
    ["인증", "Firebase Auth", "이메일/비밀번호 인증, 별도 서버 불필요"],
    ["DB", "Firebase Firestore", "실시간 동기화·NoSQL, 클라이언트 직접 연동"],
    ["차트", "Recharts / Chart.js", "통계 막대/파이 차트"],
    ["배포", "Vercel / Firebase Hosting", "정적 호스팅 + 자동 배포"],
], [Inches(2.0), Inches(4.2), Inches(5.9)], fontsize=13)

# ───────────────────────── 13. 시스템 구조 ─────────────────────────
s = add_slide(); header(s, "시스템 구조 (3-Tier · 프론트엔드 중심)", "07")
cols = [
    ("Presentation Tier", "브라우저 / Client", [
        "Next.js + React", "Tailwind 반응형 UI",
        "상태관리(로그인·그룹)", "PC 사이드바/모바일 탭"], NAVY),
    ("Application / Logic", "클라이언트 로직 + Rules", [
        "순번 진행 / 차례 복원", "통계 집계 / 랜덤 배정",
        "Security Rules 권한", "API Route(검증시만)"], RGBColor(0x1A, 0x8A, 0x7A)),
    ("Data Tier", "Firebase", [
        "Firestore", "users/groups", "chores/choreLog", "Firebase Auth"], RGBColor(0xE0, 0x8E, 0x0E)),
]
cw = Inches(3.9); gap = Inches(0.3)
x = Inches(0.55); y = Inches(1.7)
for title, sub, items, color in cols:
    fill_rect(s, x, y, cw, Inches(0.95), color)
    tf = textbox(s, x + Inches(0.1), y + Inches(0.08), cw - Inches(0.2), Inches(0.8), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; set_font(r, 16, bold=True, color=WHITE)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = sub; set_font(r, 11, color=WHITE)
    fill_rect(s, x, y + Inches(0.95), cw, Inches(3.1), LIGHT)
    btf = textbox(s, x + Inches(0.25), y + Inches(1.15), cw - Inches(0.5), Inches(2.8))
    bullets(btf, [(it, 0) for it in items], size=14)
    x += cw + gap
tf = textbox(s, Inches(0.55), Inches(6.05), Inches(12.2), Inches(1.2))
bullets(tf, [
    ("프론트엔드 중심: 클라이언트가 Firestore SDK로 직접 읽기/쓰기 (← SDK →)", 0, True),
    ("접근 제어: Firestore Security Rules로 그룹 멤버만 데이터 접근 보장", 1),
], size=14)

# ───────────────────────── 14. 데이터 모델 ─────────────────────────
s = add_slide(); header(s, "데이터 모델 (Firestore 4개 컬렉션)", "07")
add_table(s, Inches(0.6), Inches(1.5), Inches(12.1), [
    ["컬렉션", "핵심 필드"],
    ["users", "uid, name, email, groupIds[]"],
    ["groups", "name, inviteCode, ownerId, memberUids[]"],
    ["chores", "groupId, name, mode, color, rotationOrder[], currentTurnIndex, allowProxyComplete, fixedSchedule[], rules[]"],
    ["choreLog", "choreId, groupId, completedBy, completedByActual, completedAt, type, active, deactivateReason"],
], [Inches(2.0), Inches(10.1)], fontsize=13)
fill_rect(s, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.3), LIGHT)
tf = textbox(s, Inches(0.9), Inches(5.35), Inches(11.5), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
bullets(tf, [
    ("핵심 불변식", 0, True),
    ("순번 완료 시 currentTurnIndex +1(modulo) · 비활성화 시 차례 복원 · 통계는 active=true만 집계", 1),
], size=14)

# ───────────────────────── 15. 개발 일정 ─────────────────────────
s = add_slide(); header(s, "개발 일정 (3주 집중)", "08")
add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), [
    ["주차", "단계", "주요 작업", "산출물"],
    ["1주차", "기반 + 인증·그룹", "Firebase/Rules·Next.js 셋업, 회원가입/로그인, 그룹 생성·초대·전환", "레포·스키마, 인증·그룹"],
    ["2주차", "집안일 핵심 + 캘린더", "집안일 CRUD·프리셋, 순번 완료+이행, 고정제, 캘린더·팝업, 비활성화+복원", "홈·관리·캘린더"],
    ["3주차", "부가 기능 + 마감", "통계·차트, 꽝뽑기, 반응형 점검, 버그 수정, 배포, 발표", "배포본·발표"],
], [Inches(1.1), Inches(2.4), Inches(6.4), Inches(2.2)], fontsize=12)
add_table(s, Inches(0.6), Inches(4.0), Inches(12.1), [
    ["주차", "전반 (1~3일)", "후반 (4~7일)"],
    ["1주차", "프로젝트 셋업·Firestore 스키마·Rules", "인증(회원가입/로그인) + 그룹 관리"],
    ["2주차", "집안일 CRUD·순번 완료·고정제", "캘린더·상세 팝업·비활성화/복원"],
    ["3주차", "통계·차트·꽝뽑기", "반응형·배포·발표 준비"],
], [Inches(1.1), Inches(5.5), Inches(5.5)], fontsize=12)
tf = textbox(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.7))
bullets(tf, [("중간 점검(1주차 말): 인증·그룹 기능 동작 데모 + 진행률·방향성 점검", 0, True)], size=13)

# ───────────────────────── 16. 마무리 ─────────────────────────
s = add_slide()
fill_rect(s, 0, 0, SW, SH, NAVY)
tf = textbox(s, Inches(1), Inches(2.9), Inches(11.3), Inches(1.7), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "감사합니다"
set_font(r, 40, bold=True, color=WHITE)
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r = p2.add_run(); r.text = "우리집 살림 매니저 — 함께 사는 모두의 공정한 분담"
set_font(r, 18, color=RGBColor(0xBF, 0xD4, 0xEE))

prs.save(OUT)
print("SAVED", OUT, "·", len(prs.slides._sldIdLst), "slides")
